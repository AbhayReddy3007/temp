# app.py
# AI PPT Generator with improved Comparison slide support (topic detection + concise bullets)

import os
import re
import io
import tempfile
import fitz
import docx
import requests
import streamlit as st
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from datetime import date

# ---------------- CONFIG ----------------
# Replace this with an env var in production
GEMINI_API_KEY = "AIzaSyBtah4ZmuiVkSrJABE8wIjiEgunGXAbT3Q"
TEXT_MODEL_NAME = "gemini-2.0-flash"

# ---------------- LLM / GEMINI HELPERS ----------------
def call_gemini(prompt: str, timeout: int = 120) -> str:
    url = f"https://generativelanguage.googleapis.com/v1beta/models/{TEXT_MODEL_NAME}:generateContent?key={GEMINI_API_KEY}"
    payload = {"contents": [{"parts": [{"text": prompt}]}]}
    try:
        resp = requests.post(url, json=payload, timeout=timeout)
        resp.raise_for_status()
        data = resp.json()
        return data["candidates"][0]["content"]["parts"][0]["text"].strip()
    except Exception as e:
        return f"⚠️ Gemini API error: {e}"

def generate_title(summary: str) -> str:
    if not summary or not summary.strip():
        return "Presentation"
    prompt = f"""Read the following summary and create a short, clear, presentation-style title.
- Keep it under 10 words
- Do not include birth dates, long sentences, or excessive details
- Just give a clean title, like a presentation heading

Summary:
{summary}
"""
    result = call_gemini(prompt)
    return result.split("\n")[0] if result else "Presentation"

# ---------------- UTIL ----------------
def extract_slide_count(description: str, default=None):
    if not description:
        return None if default is None else max(1, default - 1)
    m = re.search(r"(\d+)\s*(slides?|sections?|pages?)", description, re.IGNORECASE)
    if m:
        total = int(m.group(1))
        return max(1, total - 1)
    return None if default is None else max(1, default - 1)

def sanitize_filename(name: str) -> str:
    return re.sub(r'[^A-Za-z0-9_.-]', '_', name).strip("_")

def clean_title_text(title: str) -> str:
    return re.sub(r"\s+", " ", title.strip()) if title else "Presentation"

def hex_to_rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.lstrip("#")
    if len(hex_color) != 6:
        hex_color = "000000"
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))

# ---------------- PARSING + OUTLINE HELPERS ----------------
def parse_points(points_text: str):
    if not points_text:
        return []

    text = points_text.replace("\r\n", "\n").strip()
    slides = []

    # 1) Slide N: split
    split_by_slide = re.split(r"\n(?=Slide\s*\d+\s*:)", text, flags=re.IGNORECASE)
    if len(split_by_slide) > 1:
        for block in split_by_slide:
            block = block.strip()
            if not block:
                continue
            lines = [ln for ln in block.splitlines() if ln.strip()]
            if not lines:
                continue
            header = lines[0]
            m = re.match(r"^\s*Slide\s*\d+\s*:\s*(.+)$", header, re.IGNORECASE)
            if m:
                title = m.group(1).strip()
                body_lines = lines[1:]
            else:
                title = lines[0].strip()
                body_lines = lines[1:]
            normalized = []
            for ln in body_lines:
                if re.match(r"^\s*[\-\u2022\*]\s+", ln):
                    normalized.append("• " + re.sub(r"^\s*[\-\u2022\*]\s*", "", ln).strip())
                elif ln.strip():
                    normalized.append(ln.strip())
            slides.append({"title": title, "description": "\n".join(normalized).strip()})
        return slides

    # 2) Double newline sections
    blocks = [b.strip() for b in re.split(r"\n\s*\n", text) if b.strip()]
    if len(blocks) > 1:
        for blk in blocks:
            lines = [l for l in blk.splitlines() if l.strip()]
            if not lines:
                continue
            title = lines[0].strip()
            rest = []
            for ln in lines[1:]:
                if re.match(r"^\s*[\-\u2022\*]\s+", ln):
                    rest.append("• " + re.sub(r"^\s*[\-\u2022\*]\s*", "", ln).strip())
                else:
                    rest.append(ln.strip())
            slides.append({"title": title, "description": "\n".join(rest).strip()})
        return slides

    # 3) short lines -> titles
    lines = [l.strip() for l in text.splitlines() if l.strip()]
    if lines:
        short_lines = [ln for ln in lines if len(ln.split()) <= 8]
        if len(short_lines) >= 3 and len(short_lines) == len(lines):
            for ln in short_lines:
                slides.append({"title": ln, "description": ""})
            return slides

    # 4) fallback single slide
    if lines:
        title = lines[0]
        body = "\n".join(lines[1:]) if len(lines) > 1 else ""
        body_lines = []
        for ln in body.splitlines():
            if re.match(r"^\s*[\-\u2022\*]\s+", ln):
                body_lines.append("• " + re.sub(r"^\s*[\-\u2022\*]\s*", "", ln).strip())
            elif ln.strip():
                body_lines.append(ln.strip())
        return [{"title": title, "description": "\n".join(body_lines).strip()}]

    return []

def generate_outline(description: str):
    if not description or not description.strip():
        return []
    num_slides = extract_slide_count(description, default=None)
    if num_slides:
        count_instruction = f"Generate {num_slides} slides."
    else:
        count_instruction = "Generate an appropriate number of slides based on the content provided."
    prompt = (
        f"Create a PowerPoint outline on: {description}\n\n"
        f"{count_instruction}\n\n"
        "IMPORTANT: Return the outline in this exact machine-friendly format and nothing else:\n\n"
        "Each slide should have a short title and 3–4 bullet points."
        "Slide 1: <Title>\n"
        "• <bullet 1>\n"
        "• <bullet 2>\n"
        "\n"
        "Slide 2: <Title>\n"
        "• <bullet 1>\n"
        "• <bullet 2>\n"
        "\n"
        "Slide 3: <Title>\n"
        "• <bullet 1>\n"
        "... and so on."
    )
    outline_text = call_gemini(prompt)
    slides = parse_points(outline_text)
    if not slides:
        retry_prompt = (
            f"Return the outline only in this exact format:\n"
            "Slide 1: <Title>\n• <bullet>\n• <bullet>\n\n"
            "Slide 2: <Title>\n• <bullet>\n\n"
            f"Topic: {description}\n{count_instruction}\n\n"
            "Return only the outline, nothing else."
        )
        outline_text2 = call_gemini(retry_prompt)
        slides = parse_points(outline_text2)
    return slides

def edit_outline_with_feedback(outline, feedback: str):
    if not outline or "slides" not in outline:
        return outline
    outline_text = "\n".join([f"Slide {i+1}: {s['title']}\n{s['description']}" for i, s in enumerate(outline['slides'])])
    prompt = (
        "Refine the outline below based on feedback.\n\n"
        "Outline:\n"
        f"{outline_text}\n\n"
        "Feedback:\n"
        f"{feedback}\n\n"
        "IMPORTANT: Return the entire updated outline in the same machine-friendly format (Slide N: ...)."
    )
    updated = call_gemini(prompt)
    updated_slides = parse_points(updated)
    if not updated_slides:
        retry = call_gemini(f"Return the updated outline only (Slide N: format). Outline:\n{outline_text}\nFeedback:\n{feedback}")
        updated_slides = parse_points(retry)
    return {"title": outline.get("title", "Presentation"), "slides": updated_slides}

# ---------------- long text summarization ----------------
def split_text(text: str, chunk_size: int = 8000, overlap: int = 300):
    chunks = []
    start = 0
    n = len(text)
    while start < n:
        end = min(start + chunk_size, n)
        chunks.append(text[start:end])
        if end == n:
            break
        start = max(0, end - overlap)
    return chunks

def summarize_long_text(full_text: str) -> str:
    """
    Produces a comprehensive, exhaustive, and structured summary of the entire document.
    - If the document is short: analyze whole document directly with a single, thorough prompt.
    - If the document is long: split into chunks, produce detailed analysis per chunk, then combine
      into one unified, exhaustive summary that preserves all important points, structure, facts,
      and nuances.
    """
    if not full_text or not full_text.strip():
        return ""

    chunks = split_text(full_text, chunk_size=8000, overlap=400)

    # If it's small, ask Gemini to produce a single exhaustive analysis
    if len(chunks) <= 1:
        prompt = f"""
Read and analyze the entire document below thoroughly. Produce a comprehensive, detailed, and exhaustive summary that preserves every important point, fact, argument, example, and nuance from the text. Do NOT oversimplify or omit material. The output should include:

1) An Executive Summary (one paragraph) that captures the overall purpose and conclusions.
2) A clear reconstruction of the document's structure with headings (e.g., Introduction, Methods/Body, Results/Arguments, Examples, Discussion, Conclusion).
3) For each section: a long, detailed section-by-section summary with important points, supporting evidence, examples, and any arguments or lines of reasoning fully preserved.
4) A consolidated list of Key Facts & Figures (as bullets), including any numbers, dates, named items, or data points.
5) Notable quotes or short excerpts (if present), labelled with approximate location.
6) Any assumptions, limitations, or open questions raised by the document.
7) A final 'Key takeaways' bullet list summarizing the most critical items.

Be exhaustive but keep the final output readable and well-structured. Document:
----------------
{full_text}
----------------
"""
        return call_gemini(prompt).strip()

    # If long, produce detailed analysis for each chunk then combine.
    partial_analyses = []
    for idx, ch in enumerate(chunks, start=1):
        prompt_chunk = f"""
You will be given CHUNK {idx} of a larger document. Carefully analyze this chunk and produce:
A) A detailed, exhaustive summary of CHUNK {idx} that preserves all important points, facts, arguments, examples, and nuance from this chunk.
B) A short heading describing what this chunk contains (e.g., "Introduction", "Methodology", "Case Study", "Analysis", "Conclusion", etc.).
C) A list of Key Facts & Figures found in this chunk (bulleted).
D) Any notable quotes or short excerpts.
E) Any open questions or references that should be cross-referenced with other chunks.

Label the output clearly as "CHUNK {idx} ANALYSIS".

Chunk content follows:
----------------
{ch}
----------------
"""
        analysis = call_gemini(prompt_chunk)
        partial_analyses.append(f"CHUNK {idx} ANALYSIS:\n{analysis.strip()}")

    combined_analyses_text = "\n\n".join(partial_analyses)

    # Combine into one final exhaustive summary
    combine_prompt = f"""
You have a set of detailed chunk analyses from a long document (listed below). Use them to produce ONE unified, coherent, and exhaustive summary of the entire original document. The final output MUST preserve every important point, fact, argument, example, and nuance found across the chunks. DO NOT INVENT new facts.

The final summary should be structured as follows:

1) Executive Summary: One concise paragraph that captures the entire document's purpose and conclusions.
2) Document Structure Reconstruction: Recreate the original document's sections and provide headings (Introduction, Body sections, Results/Arguments, Examples/Case-Studies, Discussion, Conclusion, etc.). For each reconstructed section, provide a thorough, long-form synthesis combining the chunk-level details.
3) Consolidated Key Facts & Figures: A single, deduplicated bulleted list containing all factual items (numbers, dates, names, data points) encountered in the chunks. If a fact appears in multiple chunks, include it once and list chunk locations in parentheses.
4) Important Quotes & Locations: A short list of notable quotes/excerpts and the approximate chunk number where they appear.
5) Assumptions, Limitations, and Open Questions: Combined and organized.
6) Key Takeaways: Clear bulleted summary of the most important conclusions and actionable points.

Below are the chunk analyses. Use them to reconstruct the full document and ensure no detail is lost:

----------------
{combined_analyses_text}
----------------

Now produce the final unified summary described above.
"""
    final_summary = call_gemini(combine_prompt)
    return final_summary.strip()

# ---------------- file helpers ----------------
def extract_text(path: str, filename: str) -> str:
    name = filename.lower()
    try:
        if name.endswith(".pdf"):
            doc = fitz.open(path)
            pages = []
            for page in doc:
                pages.append(page.get_text("text"))
            doc.close()
            return "\n".join(pages)
        elif name.endswith(".docx"):
            d = docx.Document(path)
            return "\n".join(p.text for p in d.paragraphs)
        elif name.endswith(".txt"):
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
    except Exception:
        return ""
    return ""

# ---------------- Comparison parsing & helpers ----------------
def extract_topics_from_feedback(feedback: str):
    """
    Attempt to extract two topics from feedback like:
     - "differentiate A and B"
     - "compare A and B"
     - "A vs B"
     - "Compare A with B"
    Returns tuple (topicA, topicB) or (None,None)
    """
    if not feedback:
        return None, None
    s = feedback.strip()
    # vs pattern
    m = re.search(r"([A-Za-z0-9 &\-\+\']{2,80})\s+vs\.?\s+([A-Za-z0-9 &\-\+\']{2,80})", s, re.IGNORECASE)
    if m:
        return m.group(1).strip(), m.group(2).strip()
    m2 = re.search(r"compare\s+(.+?)\s+(?:and|with|vs\.?)\s+(.+)", s, re.IGNORECASE)
    if m2:
        a = m2.group(1).strip().strip(",:;")
        b = m2.group(2).strip().strip(",:;")
        return a, b
    return None, None

def infer_topics_from_title(title: str):
    """
    If slide title contains 'A vs B' or 'A vs. B' or 'A vs B', extract topics.
    """
    if not title:
        return None, None
    m = re.search(r"(.+?)\s+vs\.?\s+(.+)", title, re.IGNORECASE)
    if m:
        return m.group(1).strip(), m.group(2).strip()
    return None, None

def split_long_line_into_bullets(line: str, max_words: int = 12):
    """
    Heuristic: if a line is long, split on ';' or ' — ' or commas, or sentences.
    Keep pieces with <= max_words when possible.
    Returns list of shorter strings.
    """
    if not line:
        return []
    line = line.strip()
    # If already short, return
    if len(line.split()) <= max_words:
        return [line]
    # split on semicolon or ' - ' or ' — '
    for sep in [";", "—", " - ", " – "]:
        if sep in line:
            parts = [p.strip() for p in line.split(sep) if p.strip()]
            out = []
            for p in parts:
                if len(p.split()) > max_words:
                    out.extend(split_long_line_into_bullets(p, max_words))
                else:
                    out.append(p)
            if out:
                return out
    # split on sentences
    parts = [p.strip() for p in re.split(r'(?<=[.!?])\s+', line) if p.strip()]
    if len(parts) > 1:
        out = []
        for p in parts:
            if len(p.split()) > max_words:
                out.extend(split_long_line_into_bullets(p, max_words))
            else:
                out.append(p)
        return out
    # as last resort split by commas
    if "," in line:
        parts = [p.strip() for p in line.split(",") if p.strip()]
        out = []
        for p in parts:
            if len(p.split()) > max_words:
                out.extend(split_long_line_into_bullets(p, max_words))
            else:
                out.append(p)
        return out
    # cannot split further
    return [line]

def parse_comparison_block(description: str):
    """
    Parse a description expected to contain a comparison:
    Return left_title, left_bullets(list), right_title, right_bullets(list)
    """
    left_title = ""
    right_title = ""
    left_lines = []
    right_lines = []
    if not description or not description.strip():
        return left_title or "Left", [], right_title or "Right", []

    raw = description.replace("\r\n", "\n")
    lines_raw = [ln.rstrip() for ln in raw.split("\n")]
    lines = [ln.strip() for ln in lines_raw if ln.strip()]

    mode = None
    for ln in lines:
        mleft = re.match(r"^Left\s*:\s*(.+)$", ln, re.IGNORECASE)
        mright = re.match(r"^Right\s*:\s*(.+)$", ln, re.IGNORECASE)
        if mleft:
            mode = "left"
            left_title = mleft.group(1).strip()
            continue
        if mright:
            mode = "right"
            right_title = mright.group(1).strip()
            continue
        # bullets
        if re.match(r"^[\-\u2022\*]\s+", ln) or re.match(r"^\d+\.", ln):
            cleaned = re.sub(r"^[\-\u2022\*]\s*|^\d+\.\s*", "", ln).strip()
            if mode == "right":
                right_lines.append(cleaned)
            else:
                left_lines.append(cleaned)
            continue
        # vs-line detection
        if re.search(r'\bvs\.?\b|\bversus\b|\bv\b', ln, re.IGNORECASE) and not left_title and not right_title:
            parts = re.split(r'\s+vs\.?\s+|\s+v\.?\s+|\s+versus\s+', ln, flags=re.IGNORECASE)
            if len(parts) == 2:
                left_title = parts[0].strip()
                right_title = parts[1].strip()
            continue
        # Titles if mode set
        if mode == "left" and not left_title:
            left_title = ln
            continue
        if mode == "right" and not right_title:
            right_title = ln
            continue
        # heuristics for short lines as titles
        if len(ln.split()) <= 6 and not left_title:
            left_title = ln
            continue
        if len(ln.split()) <= 6 and left_title and not right_title:
            right_title = ln
            continue
        # default append
        if mode == "right":
            right_lines.append(ln)
        else:
            left_lines.append(ln)

    # fallback: if both titles empty and description contains "A vs B" in the slide's content
    if (not left_title and not right_title):
        for ln in lines:
            m = re.search(r"(.+?)\s+vs\.?\s+(.+)", ln, re.IGNORECASE)
            if m:
                left_title = m.group(1).strip()
                right_title = m.group(2).strip()
                break

    # fallback: if right missing, try block splitting
    if (not right_lines and not right_title) and left_lines:
        blocks = [b.strip() for b in re.split(r"\n\s*\n", raw) if b.strip()]
        if len(blocks) >= 2:
            bleft = blocks[0].splitlines()
            bright = blocks[1].splitlines()
            if bleft:
                if not left_title and bleft[0].strip():
                    left_title = left_title or bleft[0].strip()
                    bleft = bleft[1:]
                for ln in bleft:
                    ln = ln.strip()
                    if not ln: continue
                    ln = re.sub(r"^[\-\u2022\*]\s*|^\d+\.\s*", "", ln).strip()
                    left_lines.append(ln)
            if bright:
                if not right_title and bright[0].strip():
                    right_title = right_title or bright[0].strip()
                    bright = bright[1:]
                for ln in bright:
                    ln = ln.strip()
                    if not ln: continue
                    ln = re.sub(r"^[\-\u2022\*]\s*|^\d+\.\s*", "", ln).strip()
                    right_lines.append(ln)
        else:
            # split half
            if len(left_lines) >= 2:
                mid = len(left_lines) // 2
                right_lines = left_lines[mid:]
                left_lines = left_lines[:mid]
                left_title = left_title or "Left"
                right_title = right_title or "Right"

    # finalize default titles
    left_title = left_title or "Left"
    right_title = right_title or "Right"

    # split long lines into multiple bullets for readability
    left_final = []
    for ln in left_lines:
        left_final.extend(split_long_line_into_bullets(ln, max_words=12))
    right_final = []
    for ln in right_lines:
        right_final.extend(split_long_line_into_bullets(ln, max_words=12))

    left_final = [l for l in left_final if l]
    right_final = [r for r in right_final if r]

    return left_title, left_final, right_title, right_final

# ---------------- PPT CREATION ----------------
def _add_image_to_slide(slide, image_bytes, left, top, width=None, height=None):
    try:
        img_stream = io.BytesIO(image_bytes)
        if width and height:
            slide.shapes.add_picture(img_stream, left, top, width=width, height=height)
        elif width:
            slide.shapes.add_picture(img_stream, left, top, width=width)
        elif height:
            slide.shapes.add_picture(img_stream, left, top, height=height)
        else:
            slide.shapes.add_picture(img_stream, left, top)
    except Exception:
        pass

def create_ppt(title, points, filename="output.pptx", title_size=30, text_size=22,
               font="Calibri", title_color="#5E2A84", text_color="#282828",
               background_color="#FFFFFF", theme="Custom",
               bg_title_path=None, bg_slide_path=None):
    """
    Create a PPTX file from `points` where each point is a dict:
      {"title": "...", "description": "..."}
    Supports per-slide formats stored in st.session_state["slide_formats"][idx] with values:
      "Full Text", "Text & Image", "Comparison"
    and per-slide uploaded images in st.session_state["slide_images"][idx] (bytes).
    Relies on helper functions available in the module:
      - hex_to_rgb(hex_color) -> RGBColor
      - _add_image_to_slide(slide, image_bytes, left, top, width=None, height=None)
      - split_long_line_into_bullets(text, max_words=12) -> list[str]
      - parse_comparison_block(description) -> (left_title, left_bullets, right_title, right_bullets)
    """
    prs = Presentation()
    title = clean_title_text(title)

    def set_bg(slide, image_path):
        if image_path and os.path.exists(image_path):
            try:
                slide.shapes.add_picture(image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
            except Exception:
                fill = slide.background.fill
                fill.solid()
                fill.fore_color.rgb = hex_to_rgb(background_color)
        else:
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = hex_to_rgb(background_color)

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_bg(slide, bg_title_path)
    tb = slide.shapes.add_textbox(Inches(1), Inches(1.6), Inches(8.5), Inches(2.2))
    tf = tb.text_frame
    tf.word_wrap = True
    try:
        tf.clear()
    except Exception:
        try:
            tf.text = ""
        except Exception:
            pass
    p = tf.add_paragraph()             
    p.text = title
    p.font.size = Pt(title_size)
    p.font.bold = True
    p.font.name = font
    p.font.color.rgb = hex_to_rgb(title_color)
    p.alignment = PP_ALIGN.LEFT

    # Content slides
    for idx, item in enumerate(points, start=1):
        key_point = clean_title_text(item.get("title", f"Slide {idx}"))
        description = item.get("description", "") or ""
        slide_format = st.session_state.get("slide_formats", {}).get(idx, st.session_state.get("slide_format", "Full Text"))

        slide = prs.slides.add_slide(prs.slide_layouts[5])
        set_bg(slide, bg_slide_path)

        # Slide title
        tb_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8.0), Inches(1.0))
        tf_title = tb_title.text_frame
        try:
            tf_title.clear()
        except Exception:
            try:
                tf_title.text = ""
            except Exception:
                pass
        p_title = tf_title.add_paragraph()
        p_title.text = key_point
        p_title.font.size = Pt(title_size)
        p_title.font.bold = True
        p_title.font.name = font
        p_title.font.color.rgb = hex_to_rgb(title_color)
        p_title.alignment = PP_ALIGN.LEFT

        # ---------- Full Text / Text & Image ----------
        if slide_format in ("Full Text", "Text & Image"):
            if description:
                if slide_format == "Text & Image":
                    tb_body = slide.shapes.add_textbox(Inches(1), Inches(1.6), Inches(5.0), Inches(4.0))
                else:
                    tb_body = slide.shapes.add_textbox(Inches(1), Inches(1.6), Inches(8.0), Inches(4.0))

                tf_body = tb_body.text_frame
                try:
                    tf_body.clear()
                except Exception:
                    try:
                        tf_body.text = ""
                    except Exception:
                        pass
                tf_body.word_wrap = True

                # Ensure each point shows with a visible dot bullet
                for line in description.splitlines():
                    if not line.strip():
                        continue
                    text_line = line.strip()
                    if text_line.startswith("•"):
                        display_text = text_line
                    elif text_line.startswith(("-", "*")):
                        display_text = "• " + text_line.lstrip("-* ").strip()
                    else:
                        display_text = "• " + text_line

                    p_body = tf_body.add_paragraph()
                    p_body.text = display_text
                    p_body.font.size = Pt(text_size)
                    p_body.font.name = font
                    p_body.font.color.rgb = hex_to_rgb(text_color)
                    p_body.level = 0

            # If Text & Image, place the image at the right side (or placeholder)
            if slide_format == "Text & Image":
                img_bytes = st.session_state.get("slide_images", {}).get(idx)
                if img_bytes:
                    left = Inches(6.0); top = Inches(1.6); width = Inches(3.0)
                    _add_image_to_slide(slide, img_bytes, left, top, width=width)
                else:
                    left = Inches(6.0); top = Inches(1.6); width = Inches(3.0); height = Inches(3.0)
                    try:
                        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
                        shape.fill.solid()
                        shape.fill.fore_color.rgb = RGBColor(240, 240, 240)
                        shape.line.color.rgb = RGBColor(200, 200, 200)
                        # some versions of python-pptx allow .text on shapes
                        try:
                            shape.text = "Image\nPlaceholder"
                        except Exception:
                            pass
                    except Exception:
                        pass

        # ---------- Comparison layout (two columns) ----------
        elif slide_format == "Comparison":
            # parse comparison content (robust parser should be available in module)
            left_title, left_bullets, right_title, right_bullets = parse_comparison_block(description or "")

            # If titles are generic, try to infer from slide title
            if (left_title.lower() == "left" or right_title.lower() == "right"):
                a, b = infer_topics_from_title(key_point)
                if a and b:
                    if left_title.lower() == "left":
                        left_title = a
                    if right_title.lower() == "right":
                        right_title = b

            # Left column title
            tb_left_title = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(4.0), Inches(0.7))
            tf_lt = tb_left_title.text_frame
            try:
                tf_lt.clear()
            except Exception:
                try:
                    tf_lt.text = ""
                except Exception:
                    pass
            p_lt = tf_lt.add_paragraph()
            p_lt.text = left_title
            p_lt.font.size = Pt(int(title_size * 0.9))
            p_lt.font.bold = True
            p_lt.font.name = font
            p_lt.font.color.rgb = hex_to_rgb(title_color)
            p_lt.alignment = PP_ALIGN.LEFT

            # Left bullets (ensure visible dot)
            tb_left = slide.shapes.add_textbox(Inches(0.6), Inches(1.9), Inches(4.0), Inches(4.0))
            tf_left = tb_left.text_frame
            try:
                tf_left.clear()
            except Exception:
                try:
                    tf_left.text = ""
                except Exception:
                    pass
            tf_left.word_wrap = True
            if left_bullets:
                for ln in left_bullets:
                    for piece in split_long_line_into_bullets(ln, max_words=12):
                        txt = piece.strip()
                        display_text = txt if txt.startswith("•") else f"• {txt}"
                        p = tf_left.add_paragraph()
                        p.text = display_text
                        p.font.size = Pt(text_size)
                        p.font.name = font
                        p.font.color.rgb = hex_to_rgb(text_color)
                        p.level = 0
            else:
                p = tf_left.add_paragraph()
                p.text = "—"
                p.font.size = Pt(text_size)
                p.font.name = font
                p.font.color.rgb = hex_to_rgb(text_color)

            # Right column title
            tb_right_title = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.0), Inches(0.7))
            tf_rt = tb_right_title.text_frame
            try:
                tf_rt.clear()
            except Exception:
                try:
                    tf_rt.text = ""
                except Exception:
                    pass
            p_rt = tf_rt.add_paragraph()
            p_rt.text = right_title
            p_rt.font.size = Pt(int(title_size * 0.9))
            p_rt.font.bold = True
            p_rt.font.name = font
            p_rt.font.color.rgb = hex_to_rgb(title_color)
            p_rt.alignment = PP_ALIGN.LEFT

            # Right bullets (ensure visible dot)
            tb_right = slide.shapes.add_textbox(Inches(5.2), Inches(1.9), Inches(4.0), Inches(4.0))
            tf_right = tb_right.text_frame
            try:
                tf_right.clear()
            except Exception:
                try:
                    tf_right.text = ""
                except Exception:
                    pass
            tf_right.word_wrap = True
            if right_bullets:
                for ln in right_bullets:
                    for piece in split_long_line_into_bullets(ln, max_words=12):
                        txt = piece.strip()
                        display_text = txt if txt.startswith("•") else f"• {txt}"
                        p = tf_right.add_paragraph()
                        p.text = display_text
                        p.font.size = Pt(text_size)
                        p.font.name = font
                        p.font.color.rgb = hex_to_rgb(text_color)
                        p.level = 0
            else:
                p = tf_right.add_paragraph()
                p.text = "—"
                p.font.size = Pt(text_size)
                p.font.name = font
                p.font.color.rgb = hex_to_rgb(text_color)

        # ---------- Footer ----------
        tb_footer = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.4))
        tf_footer = tb_footer.text_frame
        try:
            tf_footer.clear()
        except Exception:
            try:
                tf_footer.text = ""
            except Exception:
                pass
        p_footer = tf_footer.add_paragraph()
        p_footer.text = title
        p_footer.font.size = Pt(10)
        p_footer.font.name = font
        p_footer.font.color.rgb = RGBColor(120, 120, 120)
        p_footer.alignment = PP_ALIGN.RIGHT

    prs.save(filename)
    return filename


# ---------------- STREAMLIT UI ----------------
st.set_page_config(page_title="AI PPT Generator", layout="wide")
st.title("PPT Generator")

# session defaults
_defaults = {
    "messages": [],
    "doc_chat_history": [],
    "outline_chat": None,
    "summary_text": None,
    "summary_title": None,
    "title_size": 30,
    "text_size": 22,
    "font_choice": "Calibri",
    "title_color": "#5E2A84",
    "text_color": "#282828",
    "bg_color": "#FFFFFF",
    "theme": "Custom",
    "slide_format": "Full Text",
    "slide_formats": {},
    "slide_images": {},
}
for k, v in _defaults.items():
    if k not in st.session_state:
        st.session_state[k] = v

# customization
st.subheader("🎨 Customize PPT Style")
col1, col2 = st.columns(2)
with col1:
    st.session_state.title_size = st.number_input("📌 Title Font Size", min_value=12, max_value=80, value=st.session_state.title_size)
with col2:
    st.session_state.text_size = st.number_input("📝 Text Font Size", min_value=8, max_value=48, value=st.session_state.text_size)

st.session_state.font_choice = st.selectbox(
    "🔤 Font Family",
    ["Calibri", "Arial", "Times New Roman", "Verdana", "Georgia", "Helvetica", "Comic Sans MS"],
    index=0 if st.session_state.font_choice not in ["Calibri", "Arial", "Times New Roman", "Verdana", "Georgia", "Helvetica", "Comic Sans MS"] else ["Calibri", "Arial", "Times New Roman", "Verdana", "Georgia", "Helvetica", "Comic Sans MS"].index(st.session_state.font_choice)
)

col3, col4, col5 = st.columns(3)
with col3:
    st.session_state.title_color = st.color_picker("🎨 Title Color", st.session_state.title_color)
with col4:
    st.session_state.text_color = st.color_picker("📝 Text Color", st.session_state.text_color)
with col5:
    st.session_state.bg_color = st.color_picker("🌆 Background Color", st.session_state.bg_color)

st.session_state.theme = st.selectbox(
    "🎭 Select Theme",
    ["Dr.Reddys White Master", "Dr.Reddys Blue Master", "Custom"],
    index=0 if st.session_state.theme not in ["Dr.Reddys White Master", "Dr.Reddys Blue Master", "Custom"] else ["Dr.Reddys White Master", "Dr.Reddys Blue Master", "Custom"].index(st.session_state.theme)
)

st.markdown("---")

# upload doc
st.markdown("### Upload a document")
uploaded_file = st.file_uploader("Upload a document to generate a PPT from it", type=["pdf", "docx", "txt"])
if uploaded_file:
    with st.spinner("Extracting text from file..."):
        try:
            with tempfile.NamedTemporaryFile(delete=False) as tmp:
                tmp.write(uploaded_file.getvalue())
                tmp_path = tmp.name
            text = extract_text(tmp_path, uploaded_file.name)
        finally:
            try:
                os.remove(tmp_path)
            except Exception:
                pass
        if text and text.strip():
            st.session_state.summary_text = summarize_long_text(text)
            st.session_state.summary_title = generate_title(st.session_state.summary_text)
            st.success("✅ Document processed.")
        else:
            st.error("❌ Could not extract text from the uploaded file.")

# chat input

chat_prompt = st.chat_input("Enter prompt")
if chat_prompt:
    if st.session_state.summary_text:
        if any(w in chat_prompt.lower() for w in ["ppt", "slides", "presentation"]):
            with st.spinner("Generating outline from document and prompt..."):
                slides = generate_outline(st.session_state.summary_text + "\n\n" + chat_prompt)
                st.session_state.outline_chat = {"title": st.session_state.summary_title or "Presentation", "slides": slides}
                st.session_state.slide_formats = {}
                st.session_state.slide_images = {}
        else:
            st.session_state.doc_chat_history.append(("user", chat_prompt))
            reply = call_gemini(f"Answer using this document:\n{st.session_state.summary_text}\n\nQ:{chat_prompt}")
            st.session_state.doc_chat_history.append(("assistant", reply))
    else:
        st.session_state.messages.append(("user", chat_prompt))
        if any(w in chat_prompt.lower() for w in ["ppt", "slides", "presentation"]):
            with st.spinner("Generating outline..."):
                slides = generate_outline(chat_prompt)
                title = generate_title(chat_prompt)
                st.session_state.outline_chat = {"title": title or "Presentation", "slides": slides}
                st.session_state.slide_formats = {}
                st.session_state.slide_images = {}
        else:
            reply = call_gemini(chat_prompt)
            st.session_state.messages.append(("assistant", reply))
    st.rerun()

# show chat history
if st.session_state.get("messages"):
    with st.expander("Recent Chat (local)", expanded=False):
        for role, txt in st.session_state["messages"][-8:]:
            if role == "user":
                st.markdown(f"**You:** {txt}")
            else:
                st.markdown(f"**Assistant:** {txt}")

if st.session_state.get("doc_chat_history"):
    with st.expander("Document Chat (context)", expanded=False):
        for role, txt in st.session_state["doc_chat_history"][-8:]:
            if role == "user":
                st.markdown(f"**You (doc):** {txt}")
            else:
                st.markdown(f"**Assistant (doc):** {txt}")

st.markdown("---")

# Outline preview + edits
if st.session_state.outline_chat:
    outline = st.session_state.outline_chat
    st.subheader(f"📝 Preview Outline: {outline.get('title', 'Presentation')}")
    

    for idx, slide in enumerate(outline.get("slides", []), start=1):
        with st.expander(f"Slide {idx}: {slide.get('title','')}", expanded=False):
            desc = slide.get("description", "")
            if desc:
                st.markdown(desc.replace("\n", "\n\n"))
            else:
                st.write("_No content for this slide yet._")

            feedback_key = f"feedback_{idx}"
            feedback = st.text_area(f"✏️ Feedback for Slide {idx}", key=feedback_key, height=90)

            col_left, col_right = st.columns([3, 1])
            with col_left:
                current = st.session_state.get("slide_formats", {}).get(idx, st.session_state.get("slide_format", "Full Text"))
                format_key = f"format_{idx}"
                selected_format = st.selectbox(
                    f" Format for Slide {idx}",
                    ["Full Text", "Text & Image", "Comparison"],
                    index=0 if current not in ["Full Text", "Text & Image", "Comparison"] else ["Full Text", "Text & Image", "Comparison"].index(current),
                    key=format_key
                )
                st.session_state["slide_formats"][idx] = selected_format

                if selected_format == "Text & Image":
                    img_key = f"slide_image_{idx}"
                    uploaded_img = st.file_uploader(f"🖼 Upload image for Slide {idx} (optional)", type=["png", "jpg", "jpeg"], key=img_key)
                    if uploaded_img:
                        try:
                            img_bytes = uploaded_img.getvalue()
                            st.session_state["slide_images"][idx] = img_bytes
                            st.image(img_bytes, caption=f"Preview Slide {idx} image", use_column_width=True)
                        except Exception:
                            st.warning("Could not read uploaded image.")
                    else:
                        if st.session_state.get("slide_images", {}).get(idx):
                            st.image(st.session_state["slide_images"][idx], caption=f"Current Slide {idx} image (uploaded)", use_column_width=True)
                            if st.button(f"Remove image for Slide {idx}", key=f"remove_img_{idx}"):
                                st.session_state["slide_images"].pop(idx, None)
                                st.success("Image removed.")

            with col_right:
                edit_btn_key = f"edit_btn_{idx}"
                if st.button(f"💡 Edit Slide {idx}", key=edit_btn_key):
                    with st.spinner(f"Applying feedback to Slide {idx}..."):
                        # If Comparison: try to extract topics from feedback or slide title
                        topicA, topicB = extract_topics_from_feedback(feedback)
                        if not topicA or not topicB:
                            # try title
                            tA, tB = infer_topics_from_title(slide.get("title",""))
                            if tA and tB:
                                topicA, topicB = topicA or tA, topicB or tB

                        if selected_format == "Comparison":
                            # strong instruction: produce Left/Right with concise bullets
                            comp_prompt = (
                                "You are an assistant that produces a clear two-column comparison for a PowerPoint slide.\n"
                                f"Slide Title: {slide.get('title','')}\n"
                                f"Current Content: {slide.get('description','')}\n\n"
                                f"User Feedback: {feedback or '(no feedback)'}\n\n"
                                "If the user asked to compare two specific topics, ensure Left: and Right: use those topics as column titles.\n"
                                "Return output EXACTLY in this machine-friendly format and nothing else:\n\n"
                                "Left: <Topic A>\n"
                                "• <concise bullet 1 about Topic A>\n"
                                "• <concise bullet 2 about Topic A>\n"
                                "• <concise bullet 3 about Topic A>\n"
                                "• <concise bullet 4 about Topic A>\n\n"
                                "Right: <Topic B>\n"
                                "• <concise bullet 1 about Topic B>\n"
                                "• <concise bullet 2 about Topic B>\n"
                                "• <concise bullet 3 about Topic B>\n"
                                "• <concise bullet 4 about Topic B>\n\n"
                                "Bullets should be short (ideally <= 10 words). Use 'Left:' and 'Right:' labels exactly."
                            )
                            # If we inferred topics, give explicit hint
                            if topicA and topicB:
                                comp_prompt = f"Please compare '{topicA}' (Left) and '{topicB}' (Right).\n\n" + comp_prompt
                            result = call_gemini(comp_prompt)
                            # parse result: if parse_points can handle Slide headers, else store raw
                            # We'll keep result as description (parse_comparison_block will later parse)
                            if result:
                                st.session_state.outline_chat["slides"][idx - 1]["description"] = result
                                st.success(f"✅ Comparison slide {idx} updated.")
                                st.rerun()
                            else:
                                st.warning("No response from the model. Try again.")
                        else:
                            # non-comparison: reuse previous logic to request an updated slide body
                            prompt_lines = [
                                "You are an assistant that updates a PowerPoint slide.",
                                f"Slide Title: {slide.get('title','')}",
                                "Slide Content:",
                                slide.get("description", ""),
                                "",
                                "User Feedback:",
                                feedback or "(no feedback provided)",
                                "",
                                "Return only the updated bullet points or short paragraph text in a simple format."
                            ]
                            prompt = "\n".join(prompt_lines)
                            result = call_gemini(prompt)
                            parsed = parse_points(result)
                            if parsed:
                                st.session_state.outline_chat["slides"][idx - 1] = parsed[0]
                                st.success(f"✅ Slide {idx} updated.")
                                st.rerun()
                            else:
                                # fallback: store raw as description
                                if result and result.strip():
                                    st.session_state.outline_chat["slides"][idx - 1]["description"] = result
                                    st.success(f"✅ Slide {idx} updated (stored raw).")
                                    st.rerun()
                                else:
                                    st.warning("Could not parse Gemini response. Try rephrasing feedback.")

    # Outline-level controls
    st.markdown("---")
    st.subheader("Outline Controls")
    new_title = st.text_input("📌 Edit Presentation Title", value=outline.get("title", "Presentation"))
    outline_feedback = st.text_area("✏️ Feedback ", height=140)

    default_format = st.selectbox("Default Slide Format", ["Full Text", "Text & Image", "Comparison"], index=0 if st.session_state.get("slide_format") not in ["Full Text", "Text & Image", "Comparison"] else ["Full Text", "Text & Image", "Comparison"].index(st.session_state.get("slide_format", "Full Text")))
    st.session_state["slide_format"] = default_format

    col_apply, col_generate, col_clear = st.columns([1, 1, 1])
    with col_apply:
        if st.button("🔄 Apply Feedback"):
            with st.spinner("Updating outline..."):
                updated = edit_outline_with_feedback(outline, outline_feedback)
                if updated and updated.get("slides"):
                    updated["title"] = new_title.strip() if new_title else updated.get("title", "Presentation")
                    st.session_state.outline_chat = updated
                    st.success("✅ Outline updated.")
                    st.rerun()
                else:
                    st.warning("No changes returned. Try rephrasing feedback.")

    with col_generate:
        if st.button("✅ Generate PPT"):
            with st.spinner("Generating PPT..."):
                filename = f"{sanitize_filename(new_title or outline.get('title','Presentation'))}.pptx"
                if st.session_state.get("theme") == "Dr.Reddys White Master":
                    bg_title = "Screenshot 2025-10-09 163146.png"
                    bg_slide = "Screenshot 2025-10-09 163442.png"
                elif st.session_state.get("theme") == "Dr.Reddys Blue Master":
                    bg_title = "Screenshot 2025-10-09 163146.png"
                    bg_slide = "Screenshot 2025-10-09 163529.png"
                else:
                    bg_title = bg_slide = None

                create_ppt(
                    new_title or outline.get("title", "Presentation"),
                    outline.get("slides", []),
                    filename,
                    title_size=int(st.session_state.title_size),
                    text_size=int(st.session_state.text_size),
                    font=st.session_state.font_choice,
                    title_color=st.session_state.title_color,
                    text_color=st.session_state.text_color,
                    background_color=st.session_state.bg_color,
                    theme=st.session_state.theme,
                    bg_title_path=bg_title,
                    bg_slide_path=bg_slide,
                )

                try:
                    with open(filename, "rb") as f:
                        st.download_button("⬇️ Download PPT", data=f, file_name=filename, mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
                except Exception as e:
                    st.error(f"Failed to create download: {e}")



